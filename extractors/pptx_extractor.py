"""Extract products from PPTX files by extracting slide images and sending to Gemini.

Vertex AI doesn't support PPTX directly, so we:
1. Extract each slide's images from the PPTX
2. Send images to Gemini Vision for product extraction
3. Also extract text from slides for text-based product data
"""

from __future__ import annotations

import io
import logging
import re
import tempfile
from pathlib import Path

from wechat_automation.models import WeChatProduct

logger = logging.getLogger(__name__)


def extract_products_from_pptx(
    filepath: str | Path,
    source_file_id: str = "",
    vendor_id: str = "",
    vendor_name: str = "",
    max_slides: int = 50,
) -> list[WeChatProduct]:
    """Extract products from a PPTX by reading slide text + sending images to Gemini."""
    from pptx import Presentation

    path = Path(filepath)
    if not path.exists():
        return []

    try:
        prs = Presentation(str(path))
    except Exception as e:
        logger.error("Cannot open PPTX %s: %s", path.name, e)
        return []

    total_slides = len(prs.slides)
    slides_to_process = min(total_slides, max_slides)
    logger.info("Processing %s (%d slides, max %d)...", path.name, total_slides, slides_to_process)

    # Extract all text from slides
    all_text_parts: list[str] = []
    all_images: list[bytes] = []

    for i, slide in enumerate(prs.slides):
        if i >= slides_to_process:
            break

        # Extract text
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)

            # Extract table data
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        texts.append(" | ".join(row_texts))

        if texts:
            all_text_parts.append(f"--- Slide {i+1} ---\n" + "\n".join(texts))

        # Extract images from slide
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                try:
                    image_blob = shape.image.blob
                    if len(image_blob) > 5000:  # Skip tiny icons
                        all_images.append(image_blob)
                except Exception:
                    pass

    # Strategy 1: If we have meaningful text, send to Gemini as text
    combined_text = "\n\n".join(all_text_parts)
    if len(combined_text) > 200:
        products = _extract_from_text(
            combined_text,
            source_file_id=source_file_id,
            source_filename=path.name,
            vendor_id=vendor_id,
            vendor_name=vendor_name,
        )
        if products:
            logger.info("Extracted %d products from %s (text mode, %d slides)",
                        len(products), path.name, slides_to_process)
            return products

    # Strategy 2: If we have images, send the largest ones to Gemini Vision
    if all_images:
        products = _extract_from_images(
            all_images[:20],  # Max 20 images
            source_file_id=source_file_id,
            source_filename=path.name,
            vendor_id=vendor_id,
            vendor_name=vendor_name,
        )
        if products:
            logger.info("Extracted %d products from %s (image mode, %d images)",
                        len(products), path.name, len(all_images[:20]))
            return products

    logger.info("No products extracted from %s", path.name)
    return []


def _extract_from_text(
    text: str,
    source_file_id: str,
    source_filename: str,
    vendor_id: str,
    vendor_name: str,
) -> list[WeChatProduct]:
    """Send slide text to Gemini for product extraction."""
    from extractors.gemini_extractor import _init_vertex, _parse_gemini_response, MODEL_NAME, _EXTRACTION_PROMPT
    from vertexai.generative_models import GenerativeModel

    _init_vertex()
    model = GenerativeModel(MODEL_NAME)

    # Truncate if too long
    if len(text) > 100000:
        text = text[:100000]

    context = f"Vendor: {vendor_name}\n" if vendor_name else ""
    prompt = f"""{context}This is text extracted from a product presentation/catalog (PPTX).

{text}

{_EXTRACTION_PROMPT}"""

    try:
        response = model.generate_content(
            prompt,
            generation_config={"temperature": 0.1, "max_output_tokens": 65536},
        )
        if response.text:
            return _parse_gemini_response(
                response.text, source_file_id, source_filename, vendor_id, vendor_name,
            )
    except Exception as e:
        logger.error("Gemini text extraction failed for %s: %s", source_filename, e)

    return []


def _extract_from_images(
    images: list[bytes],
    source_file_id: str,
    source_filename: str,
    vendor_id: str,
    vendor_name: str,
) -> list[WeChatProduct]:
    """Send product images to Gemini Vision."""
    from extractors.gemini_extractor import _init_vertex, _parse_gemini_response, MODEL_NAME, _EXTRACTION_PROMPT
    from vertexai.generative_models import GenerativeModel, Part, Image

    _init_vertex()
    model = GenerativeModel(MODEL_NAME)

    # Build parts: images + prompt
    parts = []
    for img_bytes in images:
        # Detect mime type from magic bytes
        if img_bytes[:2] == b'\xff\xd8':
            mime = "image/jpeg"
        elif img_bytes[:4] == b'\x89PNG':
            mime = "image/png"
        else:
            mime = "image/jpeg"  # default
        parts.append(Part.from_data(data=img_bytes, mime_type=mime))

    context = f"Vendor: {vendor_name}\n" if vendor_name else ""
    parts.append(f"""{context}These are product images from a catalog presentation.

{_EXTRACTION_PROMPT}""")

    try:
        response = model.generate_content(
            parts,
            generation_config={"temperature": 0.1, "max_output_tokens": 65536},
        )
        if response.text:
            return _parse_gemini_response(
                response.text, source_file_id, source_filename, vendor_id, vendor_name,
            )
    except Exception as e:
        logger.error("Gemini image extraction failed for %s: %s", source_filename, e)

    return []
