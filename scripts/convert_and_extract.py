"""Convert DOCX/PPTX files to PDF, then extract products via Gemini.

Vertex AI doesn't support DOCX/PPTX mime types, so we convert to PDF first
using python-pptx (slide images) and python-docx (text), then send the PDF
to the standard Gemini extraction pipeline.

Usage:
    python -m scripts.convert_and_extract [--dry-run]
"""

from __future__ import annotations

import logging
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from wechat_automation import firestore_store

logger = logging.getLogger(__name__)

CONVERTIBLE_EXTS = {"docx", "doc", "pptx"}


def pptx_to_pdf(pptx_path: Path) -> Path | None:
    """Convert PPTX to PDF by rendering slides as images into a PDF."""
    from pptx import Presentation
    from pptx.util import Emu

    try:
        prs = Presentation(str(pptx_path))
    except Exception as e:
        logger.error("Cannot open PPTX %s: %s", pptx_path.name, e)
        return None

    # Extract all text and table data, write as a text-based PDF
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    tmp = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
    tmp.close()

    doc = SimpleDocTemplate(tmp.name, pagesize=A4)
    styles = getSampleStyleSheet()

    # Try to register a CJK font for Chinese text
    try:
        pdfmetrics.registerFont(TTFont('NotoSansCJK', 'C:/Windows/Fonts/msyh.ttc', subfontIndex=0))
        cjk_style = ParagraphStyle('CJK', parent=styles['Normal'], fontName='NotoSansCJK', fontSize=10)
    except Exception:
        cjk_style = styles['Normal']

    story = []
    title_style = ParagraphStyle('SlideTitle', parent=styles['Heading2'], fontName=cjk_style.fontName, fontSize=14)

    for i, slide in enumerate(prs.slides):
        story.append(Paragraph(f"Slide {i+1}", title_style))
        story.append(Spacer(1, 3 * mm))

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        # Escape XML characters
                        safe = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                        try:
                            story.append(Paragraph(safe, cjk_style))
                        except Exception:
                            pass

            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        safe = " | ".join(c.replace("&", "&amp;").replace("<", "&lt;") for c in cells)
                        try:
                            story.append(Paragraph(safe, cjk_style))
                        except Exception:
                            pass

        story.append(Spacer(1, 10 * mm))

    if not story:
        return None

    try:
        doc.build(story)
    except Exception as e:
        logger.error("PDF generation failed for %s: %s", pptx_path.name, e)
        return None

    return Path(tmp.name)


def docx_to_pdf(docx_path: Path) -> Path | None:
    """Convert DOCX to PDF by extracting text and tables."""
    from docx import Document

    try:
        doc = Document(str(docx_path))
    except Exception as e:
        logger.error("Cannot open DOCX %s: %s", docx_path.name, e)
        return None

    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    tmp = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
    tmp.close()

    pdf_doc = SimpleDocTemplate(tmp.name, pagesize=A4)
    styles = getSampleStyleSheet()

    try:
        pdfmetrics.registerFont(TTFont('NotoSansCJK', 'C:/Windows/Fonts/msyh.ttc', subfontIndex=0))
        cjk_style = ParagraphStyle('CJK', parent=styles['Normal'], fontName='NotoSansCJK', fontSize=10)
    except Exception:
        cjk_style = styles['Normal']

    story = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            safe = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            try:
                story.append(Paragraph(safe, cjk_style))
                story.append(Spacer(1, 2 * mm))
            except Exception:
                pass

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                safe = " | ".join(c.replace("&", "&amp;").replace("<", "&lt;") for c in cells)
                try:
                    story.append(Paragraph(safe, cjk_style))
                except Exception:
                    pass

    if not story:
        return None

    try:
        pdf_doc.build(story)
    except Exception as e:
        logger.error("PDF generation failed for %s: %s", docx_path.name, e)
        return None

    return Path(tmp.name)


def main() -> None:
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument("--dry-run", action="store_true")
    args = parser.parse_args()

    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s %(levelname)s %(name)s: %(message)s",
    )

    db = firestore_store._db()
    files = [d.to_dict() for d in db.collection("wechat_files").stream()]
    targets = [
        f for f in files
        if f.get("status") != "product_extracted"
        and f.get("file_extension", "") in CONVERTIBLE_EXTS
    ]

    logger.info("Found %d DOCX/PPTX files to convert and extract", len(targets))

    total_products = 0
    for i, f in enumerate(targets, 1):
        filename = f.get("filename", "?")
        ext = f.get("file_extension", "")
        source_path = f.get("source_path", "")
        file_id = f.get("file_id", "")
        vendor_id = f.get("vendor_id", "")
        vendor_name = f.get("vendor_name", "")

        local_path = Path(source_path) if source_path else None
        if not local_path or not local_path.exists():
            logger.warning("[%d/%d] File not found: %s", i, len(targets), source_path)
            continue

        logger.info("[%d/%d] Converting %s...", i, len(targets), filename[:60])

        # Convert to PDF
        pdf_path = None
        try:
            if ext == "pptx":
                pdf_path = pptx_to_pdf(local_path)
            elif ext in ("docx", "doc"):
                pdf_path = docx_to_pdf(local_path)
        except Exception as e:
            logger.error("  Conversion failed: %s", e)
            continue

        if not pdf_path or not pdf_path.exists():
            logger.info("  No content to convert")
            continue

        pdf_size = pdf_path.stat().st_size / 1024
        logger.info("  Converted to PDF (%.0fKB)", pdf_size)

        if args.dry_run:
            pdf_path.unlink(missing_ok=True)
            continue

        # Extract via Gemini
        try:
            from extractors.gemini_extractor import extract_products_gemini
            products = extract_products_gemini(
                filepath=str(pdf_path),
                source_file_id=file_id,
                vendor_id=vendor_id,
                vendor_name=vendor_name,
            )
        except Exception as e:
            logger.error("  Gemini extraction failed: %s", e)
            products = []
        finally:
            pdf_path.unlink(missing_ok=True)

        if not products:
            logger.info("  No products extracted")
            # Mark as processed anyway so we don't retry
            wf = firestore_store.get_file(file_id)
            if wf:
                wf.status = "product_extracted"
                firestore_store.upsert_file(wf)
            continue

        saved = 0
        for product in products:
            try:
                firestore_store.upsert_product(product)
                saved += 1
            except Exception:
                pass

        if saved:
            wf = firestore_store.get_file(file_id)
            if wf:
                wf.status = "product_extracted"
                firestore_store.upsert_file(wf)

        logger.info("  -> %d products extracted", saved)
        total_products += saved

    logger.info("=== Summary ===")
    logger.info("Files processed: %d", len(targets))
    logger.info("Products extracted: %d", total_products)


if __name__ == "__main__":
    main()
